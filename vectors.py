# vectors.py - Document processing and embeddings management with enhanced security and performance

import os
import logging
import hashlib
import time
import uuid
import threading
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from functools import wraps
from langchain_community.document_loaders import PyPDFLoader, UnstructuredPDFLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Qdrant
from langchain_core.documents import Document
from langchain_core.documents.base import Document

# Vector database imports
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

# Additional document loaders
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available. Install with: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install with: pip install python-pptx")

from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
load_dotenv()

# FIXED: Realistic Cost and Security limits for production use
MAX_DOCUMENT_SIZE = 50 * 1024 * 1024  # 50MB
MAX_CHUNKS_PER_DOCUMENT = 1000  # Reduced for better performance
MAX_PROCESSING_TIME = 300  # 5 minutes
ALLOWED_FILE_TYPES = {'.pdf', '.txt', '.docx', '.pptx'}
MAX_CONCURRENT_PROCESSING = 2  # Reduced to prevent resource exhaustion

# Thread-safe processing lock
processing_lock = threading.Lock()
active_sessions = {}

def retry_on_failure(max_retries=3, delay=1):
    """Decorator for retrying failed operations"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    if attempt == max_retries - 1:
                        raise e
                    logger.warning(f"Attempt {attempt + 1} failed: {e}. Retrying in {delay}s...")
                    time.sleep(delay * (attempt + 1))
            return None
        return wrapper
    return decorator

def validate_file_security(file_path: str) -> Tuple[bool, str]:
    """Enhanced security validation for uploaded files"""
    try:
        path = Path(file_path)
        
        # Check file existence
        if not path.exists():
            return False, "File does not exist"
        
        # Check file size
        file_size = path.stat().st_size
        if file_size > MAX_DOCUMENT_SIZE:
            return False, f"File too large: {file_size / (1024*1024):.1f}MB > {MAX_DOCUMENT_SIZE / (1024*1024)}MB"
        
        # Check file extension
        if path.suffix.lower() not in ALLOWED_FILE_TYPES:
            return False, f"File type not allowed: {path.suffix}"
        
        # Check for suspicious file patterns
        suspicious_patterns = ['../', '..\\', '<script', 'javascript:', 'data:', 'vbscript:', 'on load=', 'on click=']
        if any(pattern in path.name.lower() for pattern in suspicious_patterns):
            return False, "Suspicious file name detected"
        
        # Check for very small files (likely empty or corrupted)
        if file_size < 50:  # 50 bytes minimum
            return False, "File too small (likely empty or corrupted)"
        
        # Check file extension consistency (basic anti-spoofing)
        try:
            with open(file_path, 'rb') as f:
                file_header = f.read(8)
                
            # Basic file signature checks
            pdf_signatures = [b'%PDF-', b'%FDF-']
            docx_signatures = [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08']
            
            if path.suffix.lower() == '.pdf':
                if not any(file_header.startswith(sig) for sig in pdf_signatures):
                    return False, "File content doesn't match PDF format"
            elif path.suffix.lower() in ['.docx', '.pptx']:
                if not any(file_header.startswith(sig) for sig in docx_signatures):
                    return False, f"File content doesn't match {path.suffix.upper()} format"
        except:
            # If we can't read the file, it's probably corrupted
            return False, "Unable to read file - may be corrupted"
        
        return True, "File validation passed"
        
    except Exception as e:
        return False, f"Validation error: {str(e)}"

class EmbeddingsManager:
    """FIXED: Enhanced embeddings manager with better session isolation and error handling"""
    
    def __init__(
        self,
        model_name: str = "BAAI/bge-small-en",
        device: str = "cpu",
        encode_kwargs: dict = None,
        qdrant_url: str = None,
        collection_name: str = None,
        chunk_size: int = 1200,
        chunk_overlap: int = 300,
        max_chunks: int = MAX_CHUNKS_PER_DOCUMENT
    ):
        """Initialize with improved validation and session isolation"""
        
        # Generate unique session ID for this instance
        self.session_id = str(uuid.uuid4())[:8]
        
        # Validate and sanitize inputs
        self.model_name = self._sanitize_model_name(model_name)
        self.device = device if device in ["cpu", "cuda"] else "cpu"
        self.encode_kwargs = encode_kwargs or {"normalize_embeddings": True}
        self.qdrant_url = qdrant_url or os.getenv('QDRANT_URL')
        
        # IMPORTANT: Use session-specific collection
        if collection_name:
            self.collection_name = self._sanitize_collection_name(collection_name)
        else:
            # Generate unique collection name
            self.collection_name = f"doc_collection_{self.session_id}"
            
        self.api_key = os.getenv('QDRANT_API_KEY')
        
        # Validate and set processing parameters
        self.chunk_size = max(500, min(chunk_size, 2000))
        self.chunk_overlap = max(50, min(chunk_overlap, self.chunk_size // 2))
        self.max_chunks = max_chunks
        
        # Processing statistics with thread safety
        self.stats_lock = threading.Lock()
        self.stats = {
            'documents_processed': 0,
            'total_chunks_created': 0,
            'processing_time': 0,
            'errors_encountered': 0,
            'session_id': self.session_id,
            'files_processed': []  # Track processed file hashes
        }
        
        # Track this session
        with processing_lock:
            active_sessions[self.session_id] = {
                'start_time': time.time(),
                'collection_name': self.collection_name,
                'stats': self.stats
            }
        
        # Initialize embeddings model
        try:
            self.embeddings = HuggingFaceEmbeddings(
                model_name=self.model_name,
                model_kwargs={"device": self.device},
                encode_kwargs=self.encode_kwargs
            )
            logger.info(f"Session {self.session_id}: Initialized embeddings model: {self.model_name}")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to initialize embeddings: {e}")
            raise
        
        # Initialize Qdrant client
        try:
            self.qdrant_client = QdrantClient(
                url=self.qdrant_url,
                api_key=self.api_key,
                prefer_grpc=False
            )
            logger.info(f"Session {self.session_id}: Connected to Qdrant at {self.qdrant_url}")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to connect to Qdrant: {e}")
            raise
        
        # Initialize or verify collection
        self._initialize_collection()
    
    def _sanitize_model_name(self, model_name: str) -> str:
        """Sanitize model name for security"""
        allowed_models = [
            "BAAI/bge-small-en",
            "BAAI/bge-base-en",
            "sentence-transformers/all-MiniLM-L6-v2",
            "sentence-transformers/all-mpnet-base-v2"
        ]
        return model_name if model_name in allowed_models else "BAAI/bge-small-en"
    
    def _sanitize_collection_name(self, collection_name: str) -> str:
        """Sanitize collection name for security"""
        # Remove special characters, keep alphanumeric and underscores
        sanitized = ''.join(c if c.isalnum() or c == '_' else '_' for c in collection_name)
        # Limit length
        sanitized = sanitized[:64]
        # Ensure not empty
        return sanitized if sanitized else f"collection_{self.session_id}"
    
    def _initialize_collection(self) -> None:
        """Initialize or verify Qdrant collection"""
        try:
            collections = self.qdrant_client.get_collections().collections
            collection_names = [c.name for c in collections]
            
            if self.collection_name not in collection_names:
                # Create new collection
                # Fix: Use the correct vector size for the model being used
                # BAAI/bge-small-en has 384 dimensions, but let's make this more flexible
                vector_size = 384  # Default for BGE-small models
                
                # If using a different model, adjust vector size accordingly
                if "base" in self.model_name.lower():
                    vector_size = 768  # Base models typically have 768 dimensions
                elif "large" in self.model_name.lower():
                    vector_size = 1024  # Large models typically have 1024 dimensions
                
                self.qdrant_client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=VectorParams(
                        size=vector_size,
                        distance=Distance.COSINE
                    )
                )
                logger.info(f"Session {self.session_id}: Created collection: {self.collection_name} with vector size: {vector_size}")
            else:
                logger.info(f"Session {self.session_id}: Using existing collection: {self.collection_name}")
                
        except Exception as e:
            logger.error(f"Session {self.session_id}: Collection initialization failed: {e}")
            raise
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate hash of file content for duplicate detection"""
        try:
            with open(file_path, 'rb') as f:
                file_hash = hashlib.sha256(f.read()).hexdigest()
            return file_hash
        except Exception as e:
            logger.error(f"Session {self.session_id}: Hash calculation failed: {e}")
            return str(uuid.uuid4())
    
    def _load_pdf(self, file_path: str) -> List[Document]:
        """Load PDF with fallback methods"""
        loaders = [
            ("PyPDFLoader", lambda: PyPDFLoader(file_path).load()),
            ("UnstructuredPDFLoader", lambda: UnstructuredPDFLoader(file_path).load())
        ]
        
        for loader_name, loader_func in loaders:
            try:
                docs = loader_func()
                logger.info(f"Session {self.session_id}: Loaded PDF using {loader_name}")
                return docs
            except Exception as e:
                logger.warning(f"Session {self.session_id}: {loader_name} failed: {e}")
                continue
        
        raise Exception("All PDF loaders failed")
    
    def _load_txt(self, file_path: str) -> List[Document]:
        """Load text file"""
        try:
            docs = TextLoader(file_path, encoding='utf-8').load()
            logger.info(f"Session {self.session_id}: Loaded TXT file")
            return docs
        except Exception as e:
            logger.error(f"Session {self.session_id}: TXT loading failed: {e}")
            raise
    
    def _load_docx(self, file_path: str) -> List[Document]:
        """Load DOCX file"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed")
        
        try:
            doc = DocxDocument(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            content = '\n\n'.join(full_text)
            
            metadata = {
                'source': file_path,
                'file_name': Path(file_path).name,
                'file_type': 'docx'
            }
            
            logger.info(f"Session {self.session_id}: Loaded DOCX file")
            return [Document(page_content=content, metadata=metadata)]
            
        except Exception as e:
            logger.error(f"Session {self.session_id}: DOCX loading failed: {e}")
            raise
    
    def _load_pptx(self, file_path: str) -> List[Document]:
        """Load PPTX file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed")
        
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text.append(f"Slide {i}:\n" + '\n'.join(slide_text))
            
            content = '\n\n'.join(slides_text)
            
            metadata = {
                'source': file_path,
                'file_name': Path(file_path).name,
                'file_type': 'pptx',
                'total_slides': len(prs.slides)
            }
            
            logger.info(f"Session {self.session_id}: Loaded PPTX file")
            return [Document(page_content=content, metadata=metadata)]
            
        except Exception as e:
            logger.error(f"Session {self.session_id}: PPTX loading failed: {e}")
            raise
    
    @retry_on_failure(max_retries=3)
    def _load_document(self, file_path: str) -> List[Document]:
        """Load document with appropriate loader based on file type"""
        file_extension = Path(file_path).suffix.lower()
        
        loaders = {
            '.pdf': self._load_pdf,
            '.txt': self._load_txt,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx
        }
        
        loader_func = loaders.get(file_extension)
        if not loader_func:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        return loader_func(file_path)
    
    def _enhance_metadata(self, docs: List[Document], file_path: str, file_hash: str) -> List[Document]:
        """Enhance document metadata with additional information"""
        enhanced_docs = []
        file_name = Path(file_path).name
        
        for i, doc in enumerate(docs):
            # Preserve existing metadata
            metadata = doc.metadata.copy()
            
            # Add/update metadata
            metadata.update({
                'file_name': file_name,
                'file_path': file_path,
                'file_hash': file_hash,
                'chunk_index': i,
                'word_count': len(doc.page_content.split()),
                'char_count': len(doc.page_content),
                'session_id': self.session_id,
                'processed_at': time.time()
            })
            
            # Ensure page number exists
            if 'page' not in metadata:
                metadata['page'] = i + 1
            
            enhanced_docs.append(Document(
                page_content=doc.page_content,
                metadata=metadata
            ))
        
        return enhanced_docs
    
    def create_embeddings(self, file_path: str) -> str:
        """Create embeddings from document with comprehensive validation"""
        start_time = time.time()
        
        try:
            # Security validation
            is_valid, validation_message = validate_file_security(file_path)
            if not is_valid:
                raise SecurityError(f"File validation failed: {validation_message}")
            
            # Calculate file hash for duplicate detection
            file_hash = self._calculate_file_hash(file_path)
            
            # Check if file already processed
            with self.stats_lock:
                if file_hash in self.stats['files_processed']:
                    logger.info(f"Session {self.session_id}: File already processed (duplicate)")
                    return f"⚠️ File already processed: {Path(file_path).name}"
            
            logger.info(f"Session {self.session_id}: Processing file: {Path(file_path).name}")
            
            # Load document
            docs = self._load_document(file_path)
            
            if not docs:
                raise ValueError("No content extracted from document")
            
            # Enhance metadata
            docs = self._enhance_metadata(docs, file_path, file_hash)
            
            # Split documents into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            
            splits = text_splitter.split_documents(docs)
            
            # Apply chunk limit
            if len(splits) > self.max_chunks:
                logger.warning(f"Session {self.session_id}: Truncating {len(splits)} chunks to {self.max_chunks}")
                splits = splits[:self.max_chunks]
            
            if not splits:
                raise ValueError("Document splitting produced no chunks")
            
            logger.info(f"Session {self.session_id}: Created {len(splits)} valid chunks")
            
            # Create embeddings with fallback methods
            success = False
            methods = [
                ("from_documents", self._create_embeddings_method1),
                ("manual_upsert", self._create_embeddings_method2),
                ("batch_texts", self._create_embeddings_method3)
            ]
            
            for method_name, method_func in methods:
                try:
                    method_func(splits)
                    logger.info(f"Session {self.session_id}: Embeddings created using {method_name}")
                    success = True
                    break
                except Exception as e:
                    logger.warning(f"Session {self.session_id}: {method_name} failed: {e}")
                    continue
            
            if not success:
                raise ConnectionError("All embedding creation methods failed")
            
            # Update statistics thread-safely
            processing_time = time.time() - start_time
            with self.stats_lock:
                self.stats.update({
                    'documents_processed': self.stats['documents_processed'] + 1,
                    'total_chunks_created': self.stats['total_chunks_created'] + len(splits),
                    'processing_time': self.stats['processing_time'] + processing_time,
                    'files_processed': self.stats['files_processed'] + [file_hash]
                })
            
            return self._generate_success_message(file_path, splits, processing_time)
            
        except Exception as e:
            with self.stats_lock:
                self.stats['errors_encountered'] += 1
            logger.error(f"Session {self.session_id}: Embedding creation failed: {e}")
            logger.error(f"Session {self.session_id}: Full traceback:", exc_info=True)  # Add full traceback for debugging
            raise Exception(f"Embedding creation failed: {e}")

    def _create_embeddings_method1(self, splits: List[Document]) -> None:
        """Method 1: Standard from_documents"""
        try:
            # We need to ensure that the documents have valid IDs for Qdrant
            # Add UUIDs to document metadata if not present
            import uuid
            for i, doc in enumerate(splits):
                if 'id' not in doc.metadata:
                    doc.metadata['id'] = str(uuid.uuid4())
            
            Qdrant.from_documents(
                documents=splits,
                embedding=self.embeddings,
                url=self.qdrant_url,
                api_key=self.api_key,
                collection_name=self.collection_name,
                prefer_grpc=False,
            )
        except Exception as e:
            logger.error(f"Session {self.session_id}: Method 1 (from_documents) failed with detailed error: {e}")
            raise

    def _create_embeddings_method2(self, splits: List[Document]) -> None:
        """Method 2: Manual upsert with batching for better performance"""
        try:
            texts = [doc.page_content for doc in splits]
            batch_size = 10  # Reduced batch size for stability (was 25)
            
            for i in range(0, len(texts), batch_size):
                batch_texts = texts[i:i+batch_size]
                batch_docs = splits[i:i+batch_size]
                
                try:
                    embeddings_list = self.embeddings.embed_documents(batch_texts)
                    
                    points = []
                    for j, (doc, embedding) in enumerate(zip(batch_docs, embeddings_list)):
                        # Fix: Use a valid point ID format for Qdrant
                        # Generate a UUID-based point ID instead of string with underscores
                        import uuid
                        point_id = str(uuid.uuid4())
                        point = PointStruct(
                            id=point_id,
                            vector=embedding,
                            payload={
                                "page_content": doc.page_content,
                                **doc.metadata
                            }
                        )
                        points.append(point)
                    
                    self.qdrant_client.upsert(
                        collection_name=self.collection_name,
                        points=points
                    )
                    
                    logger.debug(f"Session {self.session_id}: Processed batch {i//batch_size + 1}/{(len(texts)-1)//batch_size + 1}")
                    
                except Exception as e:
                    logger.error(f"Session {self.session_id}: Batch {i//batch_size + 1} failed: {e}")
                    raise
                    
        except Exception as e:
            logger.error(f"Session {self.session_id}: Method 2 (manual_upsert) failed with detailed error: {e}")
            raise

    def _create_embeddings_method3(self, splits: List[Document]) -> None:
        """Method 3: Batch texts"""
        try:
            texts = [doc.page_content for doc in splits]
            metadatas = [doc.metadata for doc in splits]
            
            # Fix: Ensure metadatas have valid IDs for Qdrant
            import uuid
            for metadata in metadatas:
                if 'id' not in metadata:
                    metadata['id'] = str(uuid.uuid4())
            
            Qdrant.from_texts(
                texts=texts,
                embedding=self.embeddings,
                metadatas=metadatas,
                url=self.qdrant_url,
                api_key=self.api_key,
                collection_name=self.collection_name,
                prefer_grpc=False,
            )
        except Exception as e:
            logger.error(f"Session {self.session_id}: Method 3 (from_texts) failed with detailed error: {e}")
            raise

    def _generate_success_message(self, file_path: str, splits: List[Document], processing_time: float) -> str:
        """Generate detailed success message"""
        file_name = Path(file_path).name
        total_chunks = len(splits)
        avg_chunk_size = sum(len(split.page_content) for split in splits) / len(splits)
        total_words = sum(split.metadata.get('word_count', 0) for split in splits)
        
        return (
            f"✅ Successfully processed '{file_name}'!\n\n"
            f"📊 Processing Summary:\n"
            f"• Created {total_chunks} text chunks\n"
            f"• Average chunk size: {avg_chunk_size:.0f} characters\n"
            f"• Total words processed: {total_words:,}\n"
            f"• Processing time: {processing_time:.2f} seconds\n"
            f"• Collection: {self.collection_name}\n"
            f"• Session: {self.session_id}\n"
            f"• Security validation: Passed\n"
            f"• Duplicate check: Verified unique"
        )

    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics thread-safely"""
        with self.stats_lock:
            stats_copy = self.stats.copy()
        
        return {
            **stats_copy,
            "collection_name": self.collection_name,
            "max_chunks_limit": self.max_chunks,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "session_active_time": time.time() - active_sessions.get(self.session_id, {}).get('start_time', time.time()),
            "files_processed_count": len(stats_copy.get('files_processed', []))
        }

    def clear_collection(self) -> None:
        """Safely clear session-specific collection"""
        try:
            self.qdrant_client.delete_collection(self.collection_name)
            self._initialize_collection()
            
            # Reset file processing tracking
            with self.stats_lock:
                self.stats['files_processed'] = []
            
            logger.info(f"Session {self.session_id}: Collection cleared: {self.collection_name}")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to clear collection: {e}")
            raise

    def cleanup_session(self) -> None:
        """Clean up session resources"""
        try:
            # Remove from active sessions
            with processing_lock:
                if self.session_id in active_sessions:
                    del active_sessions[self.session_id]
            
            # Clean up collection
            try:
                self.qdrant_client.delete_collection(self.collection_name)
                logger.info(f"Session {self.session_id}: Cleaned up collection: {self.collection_name}")
            except Exception as e:
                logger.warning(f"Session {self.session_id}: Collection cleanup failed: {e}")
            
        except Exception as e:
            logger.error(f"Session {self.session_id}: Session cleanup failed: {e}")

    def __del__(self):
        """Cleanup when object is destroyed - BUT DON'T delete collection automatically"""
        try:
            # DON'T clean up collection here - it's still needed for chatbot queries
            # Collection cleanup should only happen on explicit session reset or app shutdown
            # Just remove from active sessions tracking
            with processing_lock:
                if self.session_id in active_sessions:
                    del active_sessions[self.session_id]
        except:
            pass

class SecurityError(Exception):
    """Custom exception for security-related errors"""
    pass

def get_active_sessions_info() -> Dict[str, Any]:
    """Get information about all active sessions"""
    with processing_lock:
        return {
            "total_active_sessions": len(active_sessions),
            "sessions": {
                session_id: {
                    "duration": time.time() - info["start_time"],
                    "collection_name": info["collection_name"],
                    "documents_processed": info["stats"]["documents_processed"]
                }
                for session_id, info in active_sessions.items()
            }
        }